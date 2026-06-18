"""
file_converter — Convert various file formats to markdown text.

Supported formats:
  .txt, .md          Plain text / markdown
  .csv               Tabular data → markdown table
  .json              JSON → fenced code block (or prettified)
  .xml               XML → fenced code block
  .html, .htm        HTML → markdown via html2text
  .pdf               PDF → page-by-page text
  .docx              Word document → plain text
  .pptx              PowerPoint → slide content
  .xlsx              Excel → sheet-by-sheet markdown tables
  .png, .jpg, .jpeg,
  .gif, .webp        Image → EXIF metadata + OCR note
  .zip               Zip archive → file listing + extracted text files
  * (other)          Raw text with unsupported-type note

Usage:
    result = convert_to_markdown("/path/to/file.pdf")
    # => {"text": "...", "filename": "file.pdf", "content_type": "...", "size": 12345}
"""

import csv
import io
import json
import os
import xml.etree.ElementTree as ET
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any


# ---------------------------------------------------------------------------
#  Image helpers
# ---------------------------------------------------------------------------
def _extract_exif(file_path: str) -> dict[str, Any]:
    """Return a dict of meaningful EXIF fields, or empty dict on failure."""
    try:
        from PIL import Image
        from PIL.ExifTags import TAGS
    except ImportError:
        return {"note": "Pillow not installed"}

    try:
        img = Image.open(file_path)
        info: dict[str, Any] = {
            "format": img.format,
            "mode": img.mode,
            "width": img.width,
            "height": img.height,
        }
        exif_data = img.getexif()
        if exif_data:
            for tag_id, value in exif_data.items():
                tag_name = TAGS.get(tag_id, tag_id)
                if isinstance(value, bytes):
                    try:
                        value = value.decode("utf-8", errors="replace")
                    except Exception:
                        value = str(value)
                # Keep only interesting fields
                if tag_name in (
                    "Make", "Model", "DateTimeOriginal", "DateTimeDigitized",
                    "ExposureTime", "FNumber", "ISOSpeedRatings",
                    "FocalLength", "Software", "Orientation", "GPSInfo",
                ):
                    if tag_name == "GPSInfo":
                        value = str(value)
                    info[tag_name] = value
        return info
    except Exception as exc:
        return {"error": str(exc)}


def _image_to_markdown(file_path: str) -> str:
    """Convert image to markdown with EXIF metadata."""
    md = [f"# Image: {os.path.basename(file_path)}", ""]
    exif = _extract_exif(file_path)
    if exif:
        md.append("## Metadata")
        for k, v in exif.items():
            md.append(f"- **{k}**: {v}")
        md.append("")
    md.append(f"![{os.path.basename(file_path)}]({file_path})")
    return "\n".join(md)


# ---------------------------------------------------------------------------
#  ZIP helper
# ---------------------------------------------------------------------------
def _zip_to_markdown(file_path: str) -> str:
    """List archive contents and extract readable text files."""
    md = [f"# ZIP Archive: {os.path.basename(file_path)}", ""]
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            md.append(f"**Total entries:** {len(zf.namelist())}")
            md.append("")
            md.append("## Contents")
            for name in zf.namelist():
                info = zf.getinfo(name)
                md.append(f"- `{name}`  ({info.file_size:,} bytes)")
            md.append("")

            # Extract small text files
            text_exts = {".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".py", ".yaml", ".yml", ".toml", ".cfg", ".ini", ".log"}
            for name in zf.namelist():
                ext = Path(name).suffix.lower()
                if ext in text_exts and zf.getinfo(name).file_size < 500_000:
                    try:
                        content = zf.read(name).decode("utf-8", errors="replace")
                        md.append(f"---")
                        md.append(f"### `{name}`")
                        md.append("")
                        md.append("```" + ext.lstrip("."))
                        md.append(content)
                        md.append("```")
                        md.append("")
                    except Exception:
                        pass
    except Exception as exc:
        md.append(f"**Error reading archive:** {exc}")
    return "\n".join(md)


# ---------------------------------------------------------------------------
#  Main conversion dispatch
# ---------------------------------------------------------------------------
def convert_to_markdown(file_path: str) -> dict[str, Any]:
    """
    Convert a file to markdown text.

    Args:
        file_path: Absolute or relative path to the file.

    Returns:
        dict with keys:
            text         — Markdown-formatted string
            filename     — Base filename
            content_type — MIME-type string
            size         — File size in bytes
    """
    path = Path(file_path)
    filename = path.name
    ext = path.suffix.lower()

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    size = path.stat().st_size

    # ── Plain text / Markdown ────────────────────────────────────────────
    if ext in (".txt", ".md"):
        text = path.read_text(encoding="utf-8", errors="replace")
        # If .md, pass through; if .txt, wrap in code fence if it looks binary-ish
        if ext == ".txt":
            if "\0" in text:
                text = f"*Binary-like content detected; showing raw decode.*\n\n```\n{text[:50000]}\n```"
        else:
            text = text
        return _result(text, filename, "text/plain", size)

    # ── CSV ──────────────────────────────────────────────────────────────
    if ext == ".csv":
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                rows = list(reader)
            if not rows:
                return _result("*Empty CSV file*", filename, "text/csv", size)

            # Build markdown table
            md = [f"# CSV: {filename}", f"**Rows:** {len(rows) - 1} (excluding header)\n"]
            header = rows[0]
            md.append("| " + " | ".join(header) + " |")
            md.append("| " + " | ".join("---" for _ in header) + " |")
            for row in rows[1:]:
                # Pad shorter rows
                row = row + [""] * (len(header) - len(row))
                md.append("| " + " | ".join(row) + " |")
            if len(rows) > 500:
                md.append(f"\n*Table truncated to {len(rows)} rows (all shown).*")
            return _result("\n".join(md), filename, "text/csv", size)
        except Exception as exc:
            return _result(f"*CSV parse error:* {exc}\n\n```\n{path.read_text(encoding='utf-8', errors='replace')[:50000]}\n```", filename, "text/csv", size)

    # ── JSON ─────────────────────────────────────────────────────────────
    if ext == ".json":
        try:
            raw = path.read_text(encoding="utf-8", errors="replace")
            parsed = json.loads(raw)
            pretty = json.dumps(parsed, indent=2, ensure_ascii=False)
            text = f"```json\n{pretty}\n```"
        except Exception:
            text = f"```\n{path.read_text(encoding='utf-8', errors='replace')[:100000]}\n```"
        return _result(text, filename, "application/json", size)

    # ── XML ──────────────────────────────────────────────────────────────
    if ext == ".xml":
        try:
            raw = path.read_text(encoding="utf-8", errors="replace")
            # Try to pretty-print via minidom
            import xml.dom.minidom
            dom = xml.dom.minidom.parseString(raw)
            pretty = dom.toprettyxml(indent="  ")
            text = f"```xml\n{pretty}\n```"
        except Exception:
            text = f"```\n{raw[:100000]}\n```"
        return _result(text, filename, "application/xml", size)

    # ── HTML ─────────────────────────────────────────────────────────────
    if ext in (".html", ".htm"):
        try:
            import html2text
            h = html2text.HTML2Text()
            h.body_width = 0  # no wrapping
            h.ignore_links = False
            h.ignore_images = False
            h.ignore_emphasis = False
            h.ignore_tables = False
            h.decode_errors = "replace"
            raw = path.read_text(encoding="utf-8", errors="replace")
            text = h.handle(raw)
        except ImportError:
            # Fallback: strip tags with BeautifulSoup
            try:
                from bs4 import BeautifulSoup
                raw = path.read_text(encoding="utf-8", errors="replace")
                soup = BeautifulSoup(raw, "html.parser")
                text = soup.get_text(separator="\n", strip=True)
                text = f"*html2text not available; using BeautifulSoup plain-text extraction.*\n\n{text}"
            except Exception:
                text = f"```html\n{raw[:100000]}\n```"
        except Exception as exc:
            raw = path.read_text(encoding="utf-8", errors="replace")
            text = f"*HTML conversion error: {exc}*\n\n```html\n{raw[:100000]}\n```"
        return _result(text, filename, "text/html", size)

    # ── PDF ──────────────────────────────────────────────────────────────
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return _result("*pypdf library not installed*", filename, "application/pdf", size)
        try:
            reader = PdfReader(file_path)
            parts = [f"# PDF: {filename}", f"**Pages:** {len(reader.pages)}\n"]
            for i, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                parts.append(f"## Page {i}")
                parts.append(page_text.strip() or "*[No extractable text on this page]*")
                parts.append("")
            return _result("\n".join(parts), filename, "application/pdf", size)
        except Exception as exc:
            return _result(f"*PDF extraction error:* {exc}", filename, "application/pdf", size)

    # ── DOCX ─────────────────────────────────────────────────────────────
    if ext == ".docx":
        try:
            from docx import Document
        except ImportError:
            return _result("*python-docx library not installed*", filename, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", size)
        try:
            doc = Document(file_path)
            parts = [f"# DOCX: {filename}", ""]
            for para in doc.paragraphs:
                parts.append(para.text)
            text = "\n".join(parts)
            # Truncate if enormous
            if len(text) > 500_000:
                text = text[:500_000] + "\n\n*[Content truncated at 500KB]*"
            return _result(text, filename, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", size)
        except Exception as exc:
            return _result(f"*DOCX extraction error:* {exc}", filename, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", size)

    # ── PPTX ─────────────────────────────────────────────────────────────
    if ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            return _result("*python-pptx library not installed*", filename, "application/vnd.openxmlformats-officedocument.presentationml.presentation", size)
        try:
            prs = Presentation(file_path)
            parts = [f"# PPTX: {filename}", f"**Slides:** {len(prs.slides)}\n"]
            for i, slide in enumerate(prs.slides, start=1):
                parts.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            parts.append(para.text)
                    if shape.has_table:
                        table = shape.table
                        tbl_md = []
                        for row_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            if row_idx == 0:
                                tbl_md.append("| " + " | ".join(cells) + " |")
                                tbl_md.append("| " + " | ".join("---" for _ in cells) + " |")
                            else:
                                tbl_md.append("| " + " | ".join(cells) + " |")
                        parts.extend(tbl_md)
                parts.append("")
            return _result("\n".join(parts), filename, "application/vnd.openxmlformats-officedocument.presentationml.presentation", size)
        except Exception as exc:
            return _result(f"*PPTX extraction error:* {exc}", filename, "application/vnd.openxmlformats-officedocument.presentationml.presentation", size)

    # ── XLSX ─────────────────────────────────────────────────────────────
    if ext == ".xlsx":
        try:
            import openpyxl
        except ImportError:
            return _result("*openpyxl library not installed*", filename, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", size)
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            parts = [f"# XLSX: {filename}", f"**Sheets:** {wb.sheetnames}\n"]
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"## {sheet_name}")
                parts.append(f"**Dimensions:** {ws.dimensions}")
                parts.append("")
                table_rows = []
                for row in ws.iter_row():
                    cells = [str(cell.value) if cell.value is not None else "" for cell in row]
                    if any(c.strip() for c in cells):
                        table_rows.append(cells)
                if table_rows:
                    header = table_rows[0]
                    parts.append("| " + " | ".join(header) + " |")
                    parts.append("| " + " | ".join("---" for _ in header) + " |")
                    for row in table_rows[1:]:
                        row = row + [""] * (len(header) - len(row))
                        parts.append("| " + " | ".join(row[:len(header)]) + " |")
                parts.append("")
            wb.close()
            return _result("\n".join(parts), filename, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", size)
        except Exception as exc:
            return _result(f"*XLSX extraction error:* {exc}", filename, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", size)

    # ── Images ───────────────────────────────────────────────────────────
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
        text = _image_to_markdown(file_path)
        mime_map = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
        }
        return _result(text, filename, mime_map.get(ext, "image/unknown"), size)

    # ── ZIP ──────────────────────────────────────────────────────────────
    if ext == ".zip":
        text = _zip_to_markdown(file_path)
        return _result(text, filename, "application/zip", size)

    # ── Unsupported types ────────────────────────────────────────────────
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        if len(raw) > 100_000:
            raw = raw[:100_000] + "\n\n*[Content truncated at 100KB]*"
    except Exception:
        raw = f"*Binary file — cannot display as text.*"
    text = (
        f"# Unsupported File Type\n\n"
        f"**Filename:** {filename}\n"
        f"**Extension:** {ext}\n"
        f"**Size:** {size:,} bytes\n\n"
        f"*This file type is not natively convertible to markdown. "
        f"Showing raw text content below (if readable).*\n\n"
        f"---\n\n"
        f"{raw}"
    )
    return _result(text, filename, "application/octet-stream", size)


def _result(text: str, filename: str, content_type: str, size: int) -> dict[str, Any]:
    """Build the standard return dict."""
    return {
        "text": text,
        "filename": filename,
        "content_type": content_type,
        "size": size,
    }


# ---------------------------------------------------------------------------
#  Convenience: convert from bytes
# ---------------------------------------------------------------------------
def convert_bytes_to_markdown(data: bytes, filename: str) -> dict[str, Any]:
    """
    Convert file content (bytes) to markdown by writing to a temp file.

    Useful for API endpoints that receive UploadFile objects.

    Args:
        data: Raw file bytes.
        filename: Original filename (used for extension detection).

    Returns:
        Same dict as convert_to_markdown().
    """
    import tempfile
    ext = Path(filename).suffix
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        result = convert_to_markdown(tmp_path)
        # Override filename to original
        result["filename"] = filename
        return result
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
