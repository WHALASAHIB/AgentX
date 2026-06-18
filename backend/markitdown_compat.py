"""
MarkItDown compatibility layer.
Provides the same API as Microsoft's markitdown but without the magika/onnxruntime
dependency that has DLL issues on this Windows system.
Uses individual file-conversion libraries that are already installed.
"""
from __future__ import annotations

import io
import os
import re
from pathlib import Path
from typing import BinaryIO, Optional, Union


class DocumentConverterResult:
    """Result of a conversion."""
    def __init__(self, title: Optional[str] = None, text_content: str = ""):
        self.title = title
        self.text_content = text_content


class MarkItDown:
    """Simple file-to-markdown converter using available libraries."""

    def __init__(self, **kwargs):
        self._requests_session = None
        pass

    def convert(self, source):
        """Convert a file to markdown. Source can be a file path or a file-like object."""
        if isinstance(source, (str, Path)):
            path = Path(source)
            suffix = path.suffix.lower()
            filename = path.name
            with open(path, "rb") as f:
                data = f.read()
        else:
            # File-like object - we need the filename
            data = source.read()
            # Try to get filename
            filename = getattr(source, 'name', 'file.bin')
            suffix = Path(filename).suffix.lower()

        text = self._convert_by_ext(data, suffix, filename)
        return DocumentConverterResult(text_content=text)

    def convert_stream(self, data, filename):
        """Convert raw bytes to markdown, given a filename for extension detection."""
        suffix = Path(filename).suffix.lower()
        text = self._convert_by_ext(data, suffix, filename)
        return DocumentConverterResult(text_content=text)

    def _convert_by_ext(self, data, suffix, filename):
        """Route to the right converter based on file extension."""
        try:
            if suffix in (".txt", ".md", ".markdown", ".csv", ".json", ".xml", ".yaml", ".yml", ".log", ".ini", ".cfg", ".conf", ".toml"):
                return self._convert_text(data)
            elif suffix in (".html", ".htm"):
                return self._convert_html(data)
            elif suffix == ".pdf":
                return self._convert_pdf(data)
            elif suffix == ".docx":
                return self._convert_docx(data)
            elif suffix == ".pptx":
                return self._convert_pptx(data)
            elif suffix in (".xlsx", ".xls"):
                return self._convert_xlsx(data, suffix)
            elif suffix in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".svg"):
                return self._convert_image(filename)
            elif suffix in (".mp3", ".wav", ".ogg", ".m4a", ".flac"):
                return self._convert_audio(filename)
            elif suffix == ".epub":
                return self._convert_epub(data)
            elif suffix == ".csv":
                return self._convert_csv(data)
            else:
                # Fallback: treat as text
                return self._convert_text(data)
        except ImportError as e:
            return f"*[Conversion failed: missing library for {suffix} files. {e}]*\n\nCould not convert {filename}."
        except Exception as e:
            return f"*[Error converting {filename}: {e}]*\n\n"

    def _convert_text(self, data):
        try:
            import charset_normalizer
            result = charset_normalizer.from_bytes(data)
            return str(result.best())
        except ImportError:
            return data.decode("utf-8", errors="replace")

    def _convert_html(self, data):
        text = data.decode("utf-8", errors="replace")
        # Remove scripts and styles
        text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
        # Simple HTML to text conversion
        text = re.sub(r'<br\s*/?>', '\n', text)
        text = re.sub(r'</p>', '\n\n', text)
        text = re.sub(r'<h([1-6])[^>]*>', lambda m: '#' * int(m.group(1)) + ' ', text)
        text = re.sub(r'</h[1-6]>', '\n', text)
        text = re.sub(r'<li[^>]*>', '- ', text)
        text = re.sub(r'</li>', '\n', text)
        text = re.sub(r'<[^>]+>', '', text)
        text = re.sub(r'&nbsp;', ' ', text)
        text = re.sub(r'&amp;', '&', text)
        text = re.sub(r'&lt;', '<', text)
        text = re.sub(r'&gt;', '>', text)
        text = re.sub(r'&quot;', '"', text)
        text = re.sub(r'&#(\d+);', lambda m: chr(int(m.group(1))), text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    def _convert_pdf(self, data):
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                pages = []
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        pages.append(t)
                return "\n\n".join(pages)
        except ImportError:
            try:
                from pdfminer.high_level import extract_text
                return extract_text(io.BytesIO(data))
            except ImportError:
                raise ImportError("pdfplumber or pdfminer")

    def _convert_docx(self, data):
        try:
            import mammoth
            result = mammoth.convert_to_markdown(io.BytesIO(data))
            return result.value
        except ImportError:
            raise ImportError("mammoth or python-docx")

    def _convert_pptx(self, data):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
            return "\n\n".join(slides)
        except ImportError:
            raise ImportError("python-pptx")

    def _convert_xlsx(self, data, suffix):
        try:
            if suffix == ".xls":
                import xlrd
                wb = xlrd.open_workbook(file_contents=data)
            else:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)

            sheets = []
            for sheet_name in wb.sheet_names:
                ws = wb[sheet_name]
                rows_text = []
                if suffix == ".xlsx":
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        rows_text.append(" | ".join(cells))
                else:
                    for row_idx in range(ws.nrows):
                        cells = [str(ws.cell_value(row_idx, c)) for c in range(ws.ncols)]
                        rows_text.append(" | ".join(cells))
                sheets.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(rows_text))
            result_text = "\n\n".join(sheets)
            if suffix == ".xlsx":
                wb.close()
            return result_text
        except ImportError:
            raise ImportError("openpyxl or xlrd")

    def _convert_image(self, filename):
        return f"![{filename}]({filename})"

    def _convert_audio(self, filename):
        return f"*[Audio file: {filename}]*"

    def _convert_epub(self, data):
        try:
            import xml.etree.ElementTree as ET
            import zipfile
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                opf_files = [f for f in zf.namelist() if f.endswith('.opf')]
                if not opf_files:
                    return "*[Unable to parse EPUB structure]*"
                opf_tree = ET.parse(io.BytesIO(zf.read(opf_files[0])))
                opf_root = opf_tree.getroot()
                items = {}
                for item in opf_root.iter('{http://www.idpf.org/2007/opf}item'):
                    item_id = item.get('id')
                    href = item.get('href')
                    if item_id and href:
                        items[item_id] = href
                spine_refs = []
                for ref in opf_root.iter('{http://www.idpf.org/2007/opf}itemref'):
                    idref = ref.get('idref')
                    if idref and idref in items:
                        spine_refs.append(items[idref])
                texts = []
                for href in spine_refs:
                    base_dir = os.path.dirname(opf_files[0])
                    full_path = os.path.join(base_dir, href) if base_dir else href
                    full_path = full_path.replace('\\', '/')
                    if full_path not in zf.namelist():
                        continue
                    content = zf.read(full_path).decode('utf-8', errors='replace')
                    texts.append(self._convert_html(content.encode()))
                return "\n\n".join(texts)
        except ImportError:
            raise ImportError("zipfile")

    def _convert_csv(self, data):
        text = data.decode("utf-8", errors="replace")
        lines = text.strip().split("\n")
        if not lines:
            return ""
        header = lines[0]
        cols = header.split(",")
        separator = "|" + "|".join(["---"] * len(cols)) + "|"
        rows = [f"|{line.replace(',', ' | ')}|" for line in lines]
        return "\n".join([rows[0], separator] + rows[1:])
